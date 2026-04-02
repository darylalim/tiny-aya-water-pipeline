from __future__ import annotations

import io

import fitz
from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches

from document import (
    _replace_paragraph_text,
    extract_segments_docx,
    extract_segments_pdf,
    extract_segments_pptx,
    extract_segments_xlsx,
    rebuild_document_docx,
    rebuild_document_pdf,
    rebuild_document_pptx,
    rebuild_document_xlsx,
    translate_document,
)


def _make_docx(paragraphs: list[str]) -> bytes:
    """Create a minimal DOCX with the given paragraph texts."""
    doc = Document()
    for text in paragraphs:
        doc.add_paragraph(text)
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _make_docx_with_table(paragraphs: list[str], table_rows: list[list[str]]) -> bytes:
    """Create a DOCX with body paragraphs and a table."""
    doc = Document()
    for text in paragraphs:
        doc.add_paragraph(text)
    table = doc.add_table(rows=len(table_rows), cols=len(table_rows[0]))
    for i, row_data in enumerate(table_rows):
        for j, cell_text in enumerate(row_data):
            table.rows[i].cells[j].text = cell_text
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _make_pptx(texts: list[str]) -> bytes:
    """Create a minimal PPTX with one slide and one textbox per text."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    for i, text in enumerate(texts):
        txBox = slide.shapes.add_textbox(
            Inches(1), Inches(0.5 + i * 1.5), Inches(5), Inches(1)
        )
        txBox.text_frame.paragraphs[0].text = text
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _make_xlsx(rows: list[list[str]]) -> bytes:
    """Create a minimal XLSX with the given rows of string values."""
    wb = Workbook()
    ws = wb.active
    for row in rows:
        ws.append(row)
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _make_pdf(texts: list[str]) -> bytes:
    """Create a minimal PDF with one text block per string."""
    doc = fitz.open()
    page = doc.new_page()
    for i, text in enumerate(texts):
        rect = fitz.Rect(72, 72 + i * 50, 400, 72 + i * 50 + 40)
        page.insert_textbox(rect, text, fontsize=12)
    buf = io.BytesIO()
    doc.save(buf)
    doc.close()
    return buf.getvalue()


# -- _replace_paragraph_text ----------------------------------------------------


def test_replace_paragraph_text_with_runs() -> None:
    """Multi-run paragraph: first run gets new text, others are cleared."""
    doc = Document()
    para = doc.add_paragraph()
    para.add_run("Hello ")
    para.add_run("World")
    _replace_paragraph_text(para, "Bonjour")
    assert para.runs[0].text == "Bonjour"
    assert para.runs[1].text == ""


def test_replace_paragraph_text_no_runs_with_text() -> None:
    """Paragraph with no runs gets a new run when replacement is non-empty."""
    doc = Document()
    para = doc.add_paragraph()
    assert len(para.runs) == 0
    _replace_paragraph_text(para, "Hello")
    assert para.text == "Hello"


def test_replace_paragraph_text_no_runs_with_whitespace() -> None:
    """Paragraph with no runs stays empty when replacement is whitespace-only."""
    doc = Document()
    para = doc.add_paragraph()
    assert len(para.runs) == 0
    _replace_paragraph_text(para, "   ")
    assert len(para.runs) == 0
    assert para.text == ""


# -- extract_segments_docx -----------------------------------------------------


def test_extract_docx_returns_paragraph_texts() -> None:
    file_bytes = _make_docx(["Hello", "World"])
    segments = extract_segments_docx(file_bytes)
    assert segments == ["Hello", "World"]


def test_extract_docx_includes_empty_paragraphs() -> None:
    file_bytes = _make_docx(["Hello", "", "World"])
    segments = extract_segments_docx(file_bytes)
    assert segments == ["Hello", "", "World"]


def test_extract_docx_includes_table_text() -> None:
    file_bytes = _make_docx_with_table(["Intro"], [["Cell A", "Cell B"]])
    segments = extract_segments_docx(file_bytes)
    assert "Intro" in segments
    assert "Cell A" in segments
    assert "Cell B" in segments


# -- rebuild_document_docx -----------------------------------------------------


def test_rebuild_docx_replaces_text() -> None:
    file_bytes = _make_docx(["Hello", "World"])
    rebuilt = rebuild_document_docx(file_bytes, ["Bonjour", "Monde"])
    segments = extract_segments_docx(rebuilt)
    assert segments == ["Bonjour", "Monde"]


def test_rebuild_docx_round_trip() -> None:
    original = ["Hello", "World"]
    file_bytes = _make_docx(original)
    segments = extract_segments_docx(file_bytes)
    rebuilt = rebuild_document_docx(file_bytes, segments)
    result = extract_segments_docx(rebuilt)
    assert result == original


def test_rebuild_docx_with_table() -> None:
    file_bytes = _make_docx_with_table(["Intro"], [["Cell A", "Cell B"]])
    segments = extract_segments_docx(file_bytes)
    translations = [s.upper() for s in segments]
    rebuilt = rebuild_document_docx(file_bytes, translations)
    result = extract_segments_docx(rebuilt)
    assert "INTRO" in result
    assert "CELL A" in result
    assert "CELL B" in result


# -- extract_segments_pptx -----------------------------------------------------


def test_extract_pptx_returns_paragraph_texts() -> None:
    file_bytes = _make_pptx(["Hello", "World"])
    segments = extract_segments_pptx(file_bytes)
    assert segments == ["Hello", "World"]


def test_extract_pptx_includes_empty_paragraphs() -> None:
    file_bytes = _make_pptx(["Hello", "", "World"])
    segments = extract_segments_pptx(file_bytes)
    assert segments == ["Hello", "", "World"]


# -- rebuild_document_pptx -----------------------------------------------------


def test_rebuild_pptx_replaces_text() -> None:
    file_bytes = _make_pptx(["Hello", "World"])
    rebuilt = rebuild_document_pptx(file_bytes, ["Bonjour", "Monde"])
    segments = extract_segments_pptx(rebuilt)
    assert segments == ["Bonjour", "Monde"]


def test_rebuild_pptx_round_trip() -> None:
    original = ["Hello", "World"]
    file_bytes = _make_pptx(original)
    segments = extract_segments_pptx(file_bytes)
    rebuilt = rebuild_document_pptx(file_bytes, segments)
    result = extract_segments_pptx(rebuilt)
    assert result == original


# -- extract_segments_xlsx -----------------------------------------------------


def test_extract_xlsx_returns_cell_texts() -> None:
    file_bytes = _make_xlsx([["Hello", "World"], ["Foo", "Bar"]])
    segments = extract_segments_xlsx(file_bytes)
    assert segments == ["Hello", "World", "Foo", "Bar"]


def test_extract_xlsx_skips_non_string_cells() -> None:
    wb = Workbook()
    ws = wb.active
    ws.append(["Hello", 42, None, "World"])
    buf = io.BytesIO()
    wb.save(buf)
    segments = extract_segments_xlsx(buf.getvalue())
    assert segments == ["Hello", "World"]


# -- rebuild_document_xlsx -----------------------------------------------------


def test_rebuild_xlsx_replaces_text() -> None:
    file_bytes = _make_xlsx([["Hello", "World"]])
    rebuilt = rebuild_document_xlsx(file_bytes, ["Bonjour", "Monde"])
    segments = extract_segments_xlsx(rebuilt)
    assert segments == ["Bonjour", "Monde"]


def test_rebuild_xlsx_round_trip() -> None:
    original = ["Hello", "World", "Foo", "Bar"]
    file_bytes = _make_xlsx([["Hello", "World"], ["Foo", "Bar"]])
    segments = extract_segments_xlsx(file_bytes)
    rebuilt = rebuild_document_xlsx(file_bytes, segments)
    result = extract_segments_xlsx(rebuilt)
    assert result == original


# -- extract_segments_pdf ------------------------------------------------------


def test_extract_pdf_returns_text_blocks() -> None:
    file_bytes = _make_pdf(["Hello", "World"])
    segments = extract_segments_pdf(file_bytes)
    assert "Hello" in segments
    assert "World" in segments


def test_extract_pdf_skips_empty_blocks() -> None:
    file_bytes = _make_pdf(["Hello", "World"])
    segments = extract_segments_pdf(file_bytes)
    assert all(s.strip() for s in segments)


# -- rebuild_document_pdf ------------------------------------------------------


def test_rebuild_pdf_inserts_translated_text() -> None:
    file_bytes = _make_pdf(["Hello", "World"])
    rebuilt = rebuild_document_pdf(file_bytes, ["Bonjour", "Monde"])
    segments = extract_segments_pdf(rebuilt)
    assert "Bonjour" in segments
    assert "Monde" in segments


def test_rebuild_pdf_preserves_segment_count() -> None:
    file_bytes = _make_pdf(["Hello", "World"])
    original_segments = extract_segments_pdf(file_bytes)
    rebuilt = rebuild_document_pdf(file_bytes, ["Bonjour", "Monde"])
    new_segments = extract_segments_pdf(rebuilt)
    assert len(new_segments) == len(original_segments)


# -- translate_document --------------------------------------------------------


def test_translate_document_dispatches_to_docx() -> None:
    file_bytes = _make_docx(["Hello", "World"])
    translated = translate_document(
        file_bytes,
        "test.docx",
        translate_fn=lambda text: text.upper(),
    )
    segments = extract_segments_docx(translated)
    assert segments == ["HELLO", "WORLD"]


def test_translate_document_skips_empty_segments() -> None:
    file_bytes = _make_docx(["Hello", "", "World"])
    calls: list[str] = []

    def mock_translate(text: str) -> str:
        calls.append(text)
        return text.upper()

    translated = translate_document(
        file_bytes, "test.docx", translate_fn=mock_translate
    )
    segments = extract_segments_docx(translated)
    assert segments == ["HELLO", "", "WORLD"]
    assert calls == ["Hello", "World"]


def test_translate_document_dispatches_to_pptx() -> None:
    file_bytes = _make_pptx(["Hello", "World"])
    translated = translate_document(
        file_bytes,
        "slides.pptx",
        translate_fn=lambda text: text.upper(),
    )
    segments = extract_segments_pptx(translated)
    assert segments == ["HELLO", "WORLD"]


def test_translate_document_dispatches_to_xlsx() -> None:
    file_bytes = _make_xlsx([["Hello", "World"]])
    translated = translate_document(
        file_bytes,
        "data.xlsx",
        translate_fn=lambda text: text.upper(),
    )
    segments = extract_segments_xlsx(translated)
    assert segments == ["HELLO", "WORLD"]


def test_translate_document_dispatches_to_pdf() -> None:
    file_bytes = _make_pdf(["Hello"])
    translated = translate_document(
        file_bytes,
        "report.pdf",
        translate_fn=lambda text: text.upper(),
    )
    segments = extract_segments_pdf(translated)
    assert any("HELLO" in s for s in segments)


def test_translate_document_preserves_whitespace_only_segments() -> None:
    file_bytes = _make_docx(["Hello", "   ", "World"])
    calls: list[str] = []

    def mock_translate(text: str) -> str:
        calls.append(text)
        return text.upper()

    translated = translate_document(
        file_bytes, "test.docx", translate_fn=mock_translate
    )
    segments = extract_segments_docx(translated)
    assert segments == ["HELLO", "   ", "WORLD"]
    assert calls == ["Hello", "World"]


def test_translate_document_unsupported_format() -> None:
    import pytest

    with pytest.raises(ValueError, match="Unsupported file format"):
        translate_document(b"data", "test.txt", translate_fn=lambda t: t)

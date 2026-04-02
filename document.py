from __future__ import annotations

import io
from pathlib import Path
from typing import Any, Callable

# -- Helpers -------------------------------------------------------------------


def _replace_paragraph_text(para: Any, text: str) -> None:
    """Replace all text in a paragraph, preserving the first run's formatting."""
    if para.runs:
        para.runs[0].text = text
        for run in para.runs[1:]:
            run.text = ""
    elif text.strip():
        para.add_run().text = text


# -- DOCX ---------------------------------------------------------------------


def _iter_docx_paragraphs(doc: Any) -> list[Any]:
    """Return all paragraphs in document body and tables, deduped."""
    seen: set[int] = set()
    paragraphs: list[Any] = []
    for para in doc.paragraphs:
        pid = id(para._element)
        if pid not in seen:
            seen.add(pid)
            paragraphs.append(para)
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for para in cell.paragraphs:
                    pid = id(para._element)
                    if pid not in seen:
                        seen.add(pid)
                        paragraphs.append(para)
    return paragraphs


def extract_segments_docx(file_bytes: bytes) -> list[str]:
    """Extract translatable text segments from a DOCX file."""
    from docx import Document

    doc = Document(io.BytesIO(file_bytes))
    return [para.text for para in _iter_docx_paragraphs(doc)]


def rebuild_document_docx(file_bytes: bytes, translations: list[str]) -> bytes:
    """Rebuild a DOCX file with translated text replacing original segments."""
    from docx import Document

    doc = Document(io.BytesIO(file_bytes))
    for i, para in enumerate(_iter_docx_paragraphs(doc)):
        if i >= len(translations):
            break
        _replace_paragraph_text(para, translations[i])
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


# -- PPTX ---------------------------------------------------------------------


def _iter_pptx_paragraphs(prs: Any) -> list[Any]:
    """Return all paragraphs across all slides, shapes, and tables."""
    paragraphs: list[Any] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    paragraphs.append(para)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for para in cell.text_frame.paragraphs:
                            paragraphs.append(para)
    return paragraphs


def extract_segments_pptx(file_bytes: bytes) -> list[str]:
    """Extract translatable text segments from a PPTX file."""
    from pptx import Presentation

    prs = Presentation(io.BytesIO(file_bytes))
    return [para.text for para in _iter_pptx_paragraphs(prs)]


def rebuild_document_pptx(file_bytes: bytes, translations: list[str]) -> bytes:
    """Rebuild a PPTX file with translated text replacing original segments."""
    from pptx import Presentation

    prs = Presentation(io.BytesIO(file_bytes))
    for i, para in enumerate(_iter_pptx_paragraphs(prs)):
        if i >= len(translations):
            break
        _replace_paragraph_text(para, translations[i])
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# -- XLSX ---------------------------------------------------------------------


def extract_segments_xlsx(file_bytes: bytes) -> list[str]:
    """Extract translatable text segments from an XLSX file."""
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(file_bytes))
    segments: list[str] = []
    for ws in wb.worksheets:
        for row in ws.iter_rows():
            for cell in row:
                if isinstance(cell.value, str):
                    segments.append(cell.value)
    return segments


def rebuild_document_xlsx(file_bytes: bytes, translations: list[str]) -> bytes:
    """Rebuild an XLSX file with translated text replacing original segments."""
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(file_bytes))
    idx = 0
    for ws in wb.worksheets:
        for row in ws.iter_rows():
            for cell in row:
                if isinstance(cell.value, str):
                    if idx < len(translations):
                        cell.value = translations[idx]
                        idx += 1
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


# -- PDF ----------------------------------------------------------------------


def extract_segments_pdf(file_bytes: bytes) -> list[str]:
    """Extract translatable text segments from a PDF file."""
    import fitz

    doc = fitz.open(stream=file_bytes, filetype="pdf")
    try:
        segments: list[str] = []
        for page in doc:
            blocks = page.get_text("dict")["blocks"]
            for block in blocks:
                if "lines" not in block:
                    continue
                text = ""
                for line in block["lines"]:
                    for span in line["spans"]:
                        text += span["text"]
                text = text.strip()
                if text:
                    segments.append(text)
        return segments
    finally:
        doc.close()


def rebuild_document_pdf(file_bytes: bytes, translations: list[str]) -> bytes:
    """Rebuild a PDF with translated text replacing original text blocks.

    Best-effort: complex layouts may not reconstruct perfectly.
    """
    import fitz

    doc = fitz.open(stream=file_bytes, filetype="pdf")
    try:
        idx = 0
        for page in doc:
            blocks = page.get_text("dict")["blocks"]
            insertions: list[tuple[Any, str, float]] = []
            for block in blocks:
                if "lines" not in block:
                    continue
                text = ""
                for line in block["lines"]:
                    for span in line["spans"]:
                        text += span["text"]
                text = text.strip()
                if not text:
                    continue
                rect = fitz.Rect(block["bbox"])
                fontsize: float = block["lines"][0]["spans"][0]["size"]
                page.add_redact_annot(rect)
                if idx < len(translations):
                    insertions.append((rect, translations[idx], fontsize))
                    idx += 1
            page.apply_redactions()
            page_width = page.rect.width
            for rect, trans_text, fontsize in insertions:
                # Expand rect to page width and add vertical room — the
                # extracted bbox is tight around glyphs so translated text
                # (often longer) would overflow the original bounds.
                expanded = fitz.Rect(
                    rect.x0,
                    rect.y0,
                    page_width - 36,
                    rect.y1 + fontsize * 2,
                )
                page.insert_textbox(expanded, trans_text, fontsize=fontsize)
        buf = io.BytesIO()
        doc.save(buf)
        return buf.getvalue()
    finally:
        doc.close()


# -- Coordinator ---------------------------------------------------------------

_ExtractFn = Callable[[bytes], list[str]]
_RebuildFn = Callable[[bytes, list[str]], bytes]

HANDLERS: dict[str, tuple[_ExtractFn, _RebuildFn]] = {
    ".docx": (extract_segments_docx, rebuild_document_docx),
    ".pdf": (extract_segments_pdf, rebuild_document_pdf),
    ".pptx": (extract_segments_pptx, rebuild_document_pptx),
    ".xlsx": (extract_segments_xlsx, rebuild_document_xlsx),
}


def translate_document(
    file_bytes: bytes,
    filename: str,
    translate_fn: Callable[[str], str],
) -> bytes:
    """Translate a document file, preserving formatting.

    Args:
        file_bytes: Raw file content.
        filename: Original filename (used to determine format from extension).
        translate_fn: A callable that takes a text string and returns its translation.

    Returns:
        The translated document as bytes in the same format.
    """
    ext = Path(filename).suffix.lower()
    if ext not in HANDLERS:
        msg = f"Unsupported file format: {ext}"
        raise ValueError(msg)
    extract_fn, rebuild_fn = HANDLERS[ext]
    segments = extract_fn(file_bytes)
    translations = [translate_fn(seg) if seg.strip() else seg for seg in segments]
    return rebuild_fn(file_bytes, translations)
